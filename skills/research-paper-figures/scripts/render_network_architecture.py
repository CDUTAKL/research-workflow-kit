"""Render neural-network architecture diagrams from a JSON structure spec.

This script is intentionally separate from result plotting. It renders model
architecture figures using reusable vector-style components: feature-map stacks,
operation tags, residual paths, repeated-block badges, insets, and classifier heads.
"""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

import matplotlib as mpl
import matplotlib.pyplot as plt
from matplotlib.patches import Circle, FancyArrowPatch, FancyBboxPatch, Rectangle


INK = "#17212B"
MUTED = "#5E6A75"
GRID = "#CFD7DF"
PAPER = "#FFFFFF"
SKIP = "#C28D3A"
BLUE_ACCENT = "#2E6F9E"
FUSION = "#8B5FBF"


def setup_matplotlib() -> None:
    mpl.rcParams.update(
        {
            "font.family": "sans-serif",
            "font.sans-serif": ["Arial", "DejaVu Sans", "Liberation Sans", "sans-serif"],
            "svg.fonttype": "none",
            "pdf.fonttype": 42,
            "font.size": 7,
        }
    )


def load_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def rgba(color: str, alpha: float):
    r, g, b, _ = mpl.colors.to_rgba(color)
    return (r, g, b, alpha)


def add_text(ax, x, y, text, size=6, color=INK, weight="normal", ha="center", va="center", **kwargs):
    ax.text(
        x,
        y,
        text,
        fontsize=size,
        color=color,
        weight=weight,
        ha=ha,
        va=va,
        linespacing=1.12,
        **kwargs,
    )


def add_arrow(ax, start, end, color=INK, lw=0.85, rad=0.0, scale=9):
    ax.add_patch(
        FancyArrowPatch(
            start,
            end,
            arrowstyle="-|>",
            mutation_scale=scale,
            linewidth=lw,
            color=color,
            shrinkA=2,
            shrinkB=2,
            connectionstyle=f"arc3,rad={rad}",
            zorder=20,
        )
    )


def feature_stack(ax, x, y, w, h, n, color, title, shape, badge=None):
    dx, dy = 0.045, 0.035
    max_planes = max(1, min(int(n), 12))
    for i in range(max_planes):
        ax.add_patch(
            Rectangle(
                (x + i * dx, y + i * dy),
                w,
                h,
                linewidth=0.48,
                edgecolor=INK,
                facecolor=rgba(color, 0.33 + i * 0.035),
                zorder=2 + i,
            )
        )

    cx = x + w / 2 + max_planes * dx / 2
    cy = y + h / 2 + max_planes * dy / 2
    ax.plot([x, x + max_planes * dx], [y + h, y + h + max_planes * dy], color=INK, lw=0.4)
    ax.plot([x + w, x + w + max_planes * dx], [y + h, y + h + max_planes * dy], color=INK, lw=0.4)
    ax.plot([x + w, x + w + max_planes * dx], [y, y + max_planes * dy], color=INK, lw=0.4)
    add_text(ax, cx, y + h + max_planes * dy + 0.18, title, size=7.0, weight="bold")
    add_text(ax, cx, y - 0.15, shape, size=5.8, color=MUTED)

    if badge:
        ax.add_patch(
            FancyBboxPatch(
                (cx - 0.31, y + h / 2 - 0.16),
                0.62,
                0.32,
                boxstyle="round,pad=0.012,rounding_size=0.08",
                facecolor=PAPER,
                edgecolor=INK,
                linewidth=0.55,
                zorder=30,
            )
        )
        add_text(ax, cx, y + h / 2, badge, size=5.4, weight="bold", zorder=31)
    return {"cx": cx, "cy": cy, "left": x, "right": x + w + max_planes * dx, "top": y + h + max_planes * dy, "bottom": y}


def operation_tag(ax, x, y, text, color, width=0.86):
    ax.add_patch(
        FancyBboxPatch(
            (x - width / 2, y - 0.17),
            width,
            0.34,
            boxstyle="round,pad=0.012,rounding_size=0.055",
            facecolor=rgba(color, 0.22),
            edgecolor=color,
            linewidth=0.7,
            zorder=10,
        )
    )
    add_text(ax, x, y, text, size=5.4, zorder=11)


def add_plus(ax, x, y):
    ax.add_patch(Circle((x, y), 0.12, facecolor=PAPER, edgecolor=INK, linewidth=0.6, zorder=25))
    ax.plot([x - 0.06, x + 0.06], [y, y], color=INK, lw=0.8, zorder=26)
    ax.plot([x, x], [y - 0.06, y + 0.06], color=INK, lw=0.8, zorder=26)


def add_fusion_node(ax, x, y, label="+"):
    ax.add_patch(Circle((x, y), 0.13, facecolor="#FAF7FF", edgecolor=FUSION, linewidth=0.7, zorder=25))
    add_text(ax, x, y, label, size=6.0, color=FUSION, weight="bold", zorder=26)


def node_lookup(input_center: dict[str, Any], stage_centers: list[dict[str, Any]], head_x: float, spec: dict[str, Any]) -> dict[str, dict[str, float]]:
    lookup: dict[str, dict[str, float]] = {
        "Input": input_center,
        spec.get("input", {}).get("name", "Input"): input_center,
        "GAP": {"cx": head_x, "cy": input_center["cy"], "left": head_x - 0.26, "right": head_x + 0.26},
    }
    for center in stage_centers:
        stage = center.get("spec", {})
        lookup[stage.get("name", "")] = center
        if stage.get("id"):
            lookup[stage["id"]] = center
    head = spec.get("head", [])
    for i, part in enumerate(head):
        name = part.get("name")
        if name:
            lookup[name] = {"cx": head_x + i * 1.08, "cy": input_center["cy"], "left": head_x + i * 1.08 - 0.25, "right": head_x + i * 1.08 + 0.25}
    return {k: v for k, v in lookup.items() if k}


def draw_declared_connections(ax, connections: list[dict[str, Any]], lookup: dict[str, dict[str, float]]) -> list[str]:
    """Draw optional explicit topology connections from the JSON spec.

    The renderer remains layout-light on purpose: the JSON spec is the source of
    truth, while this function adds visible residual, skip, fusion, attention, or
    concat paths when the spec declares them.
    """
    unresolved: list[str] = []
    for conn in connections:
        source = conn.get("from")
        target = conn.get("to")
        if source not in lookup or target not in lookup:
            unresolved.append(f"{source}->{target}")
            continue
        kind = conn.get("type", "skip")
        label = conn.get("label", kind)
        source_node = lookup[source]
        target_node = lookup[target]
        y = max(source_node.get("cy", 3.25), target_node.get("cy", 3.25)) + float(conn.get("offset", 1.0))
        color = {
            "residual": SKIP,
            "skip": SKIP,
            "concat": FUSION,
            "fusion": FUSION,
            "attention": "#4A7A8C",
        }.get(kind, MUTED)
        ax.plot(
            [source_node["cx"], source_node["cx"], target_node["cx"], target_node["cx"]],
            [source_node["cy"] + 0.28, y, y, target_node["cy"] + 0.28],
            color=color,
            lw=0.82,
            zorder=18,
        )
        if kind in {"concat", "fusion", "attention"}:
            add_fusion_node(ax, target_node["cx"], target_node["cy"] + 0.28, "+" if kind != "attention" else "g")
        add_text(ax, (source_node["cx"] + target_node["cx"]) / 2, y + 0.11, label, size=4.8, color=color)
    return unresolved


def draw_basicblock(ax, x, y, title="Residual BasicBlock", subtitle="repeated block"):
    add_text(ax, x, y + 1.2, title, size=7.1, weight="bold", ha="left")
    add_text(ax, x, y + 0.98, subtitle, size=5.5, color=MUTED, ha="left")
    boxes = [
        (x + 0.0, "x", 0.34),
        (x + 0.72, "3x3 conv\nBN, ReLU", 0.84),
        (x + 1.86, "3x3 conv\nBN", 0.74),
        (x + 2.9, "+", 0.32),
        (x + 3.52, "ReLU", 0.5),
    ]
    last_right = None
    for bx, label, bw in boxes:
        ax.add_patch(
            FancyBboxPatch(
                (bx, y + 0.42),
                bw,
                0.36,
                boxstyle="round,pad=0.012,rounding_size=0.04",
                facecolor=PAPER if label in {"x", "+", "ReLU"} else "#F5F7FA",
                edgecolor=INK,
                linewidth=0.55,
                zorder=10,
            )
        )
        add_text(
            ax,
            bx + bw / 2,
            y + 0.6,
            label,
            size=5.4 if label != "+" else 7.0,
            weight="bold" if label == "+" else "normal",
            zorder=11,
        )
        if last_right is not None:
            add_arrow(ax, (last_right, y + 0.6), (bx, y + 0.6), scale=6.5, lw=0.65)
        last_right = bx + bw
    ax.plot([x + 0.17, x + 0.17, x + 3.06, x + 3.06], [y + 0.42, y + 0.18, y + 0.18, y + 0.42], color=SKIP, lw=0.8)
    add_text(ax, x + 1.62, y + 0.06, "identity / projection shortcut", size=5.0, color=SKIP)


def stage_dimensions(index: int, count: int):
    # Gradually shrink spatial footprint while increasing stack depth.
    w = max(0.38, 0.52 - index * 0.026)
    h = max(0.28, 0.72 - index * 0.07)
    return w, h


def render(spec: dict[str, Any], preset: dict[str, Any], out: Path, formats: list[str], qa_report: bool, pptx_mode: str) -> None:
    setup_matplotlib()
    fig_size = preset.get("figure_size", [13.4, 6.0])
    dpi = int(preset.get("dpi", 450))
    palette = preset["palette"]
    fig, ax = plt.subplots(figsize=fig_size)
    ax.set_xlim(0, 13.4)
    ax.set_ylim(0, 6.0)
    ax.axis("off")
    fig.patch.set_facecolor(PAPER)

    audit = spec.get("audit", {})
    layout = spec.get("layout", "linear")
    add_text(ax, 0.42, 5.63, spec.get("title", spec.get("model", "Network architecture")), size=12.0, weight="bold", ha="left", va="top")
    add_text(ax, 0.42, 5.33, spec.get("subtitle", ""), size=6.7, color=MUTED, ha="left", va="top")
    add_text(ax, 12.75, 5.55, f"layout: {layout}", size=5.2, color=MUTED, ha="right", va="top")

    base_y = 3.25
    input_spec = spec["input"]
    input_center = feature_stack(
        ax,
        0.48,
        base_y - 0.44,
        0.44,
        0.88,
        input_spec.get("planes", 3),
        palette.get(input_spec.get("color_role", "input"), "#E8EDF2"),
        input_spec.get("name", "Input"),
        input_spec.get("shape", ""),
    )

    stages = spec.get("stages", [])
    stage_xs = [1.92, 3.38, 4.94, 6.58, 8.28, 9.82]
    stage_centers = []

    # Insert first operation before the first stage when provided.
    operations = spec.get("operations", [])
    if operations:
        op = operations[0]
        operation_tag(ax, 1.42, base_y, op["label"], palette.get(op.get("color_role", "stem"), "#7EAFD2"), 0.78)
        add_arrow(ax, (input_center["right"] + 0.1, base_y), (1.1, base_y), scale=6)

    for i, stage_spec in enumerate(stages):
        x = stage_xs[min(i, len(stage_xs) - 1)]
        w, h = stage_dimensions(i, len(stages))
        y = base_y - h / 2 + min(i, 4) * 0.03
        center = feature_stack(
            ax,
            x,
            y,
            w,
            h,
            stage_spec.get("planes", 6),
            palette.get(stage_spec.get("color_role", f"stage{i + 1}"), "#BFD7EA"),
            stage_spec["name"],
            stage_spec.get("shape", ""),
            stage_spec.get("repeat"),
        )
        center["spec"] = stage_spec
        stage_centers.append(center)

        if i == 0 and len(operations) > 1:
            op = operations[1]
            operation_tag(ax, (stage_centers[0]["left"] + input_center["right"]) / 2 + 0.35, base_y, op["label"], palette.get(op.get("color_role", "stem"), "#7EAFD2"), 0.82)

    # Main arrows.
    previous_right = input_center["right"] + 0.55
    for i, center in enumerate(stage_centers):
        start = (previous_right, base_y)
        end = (center["left"] - 0.08, base_y)
        color = BLUE_ACCENT if center["spec"].get("downsample") else INK
        add_arrow(ax, start, end, color=color, scale=8)
        if center["spec"].get("downsample"):
            add_text(ax, (start[0] + end[0]) / 2, base_y + 0.22, center["spec"]["downsample"], size=5.2, color=color)
        previous_right = center["right"] + 0.12

    # Residual shortcuts for repeated block stages.
    for i, center in enumerate(stage_centers):
        if center["spec"].get("block", "").lower().find("basic") >= 0:
            sx = center["left"] + 0.12
            ex = center["right"] - 0.1
            yoff = base_y + 0.88 + i * 0.11
            ax.plot([sx, sx, ex, ex], [base_y + 0.24, yoff, yoff, base_y + 0.24], color=SKIP, lw=0.75)
            add_plus(ax, ex, base_y + 0.24)
            add_text(ax, (sx + ex) / 2, yoff + 0.1, "projection skip" if center["spec"].get("downsample") else "skip", size=4.8, color=SKIP)

    # Classifier head.
    head_x = 10.2
    if stage_centers:
        add_arrow(ax, (stage_centers[-1]["right"] + 0.1, base_y), (head_x - 0.34, base_y), scale=8)
    head = spec.get("head", [])
    if head:
        ax.add_patch(Circle((head_x, base_y), 0.26, facecolor="#F7F2E8", edgecolor=INK, linewidth=0.6, zorder=10))
        add_text(ax, head_x, base_y, head[0].get("name", "GAP"), size=6.0, weight="bold", zorder=11)
        add_text(ax, head_x, base_y - 0.47, head[0].get("shape", ""), size=5.3, color=MUTED)
        if len(head) > 1:
            add_arrow(ax, (head_x + 0.3, base_y), (head_x + 0.72, base_y), scale=8)
            ax.add_patch(
                FancyBboxPatch(
                    (head_x + 0.76, base_y - 0.28),
                    0.86,
                    0.56,
                    boxstyle="round,pad=0.016,rounding_size=0.04",
                    facecolor=rgba(palette.get("head", "#D8837B"), 0.45),
                    edgecolor=INK,
                    linewidth=0.6,
                    zorder=10,
                )
            )
            add_text(ax, head_x + 1.19, base_y, f"{head[1].get('name', 'FC')}\n{head[1].get('shape', '')}", size=5.8, zorder=11)
        if len(head) > 2:
            add_arrow(ax, (head_x + 1.64, base_y), (head_x + 2.1, base_y), scale=8)
            add_text(ax, head_x + 2.32, base_y, head[2].get("name", "output"), size=6.0, weight="bold")
            add_text(ax, head_x + 2.32, base_y - 0.26, head[2].get("shape", ""), size=5.0, color=MUTED)

    explicit_unresolved = draw_declared_connections(ax, spec.get("connections", []), node_lookup(input_center, stage_centers, head_x, spec))

    # Shape axis.
    ax.plot([0.58, 12.5], [2.22, 2.22], color=GRID, lw=0.6)
    markers = [(input_center["cx"], input_spec.get("shape", "").split("x")[0].strip())]
    markers.extend((c["cx"], c["spec"].get("shape", "").split("x")[0].strip()) for c in stage_centers)
    markers.append((head_x, "1"))
    for x, lab in markers:
        ax.plot([x, x], [2.34, 2.10], color=GRID, lw=0.6)
        add_text(ax, x, 1.96, lab, size=5.2, color=MUTED)
    add_text(ax, 0.48, 2.48, "spatial size", size=5.6, color=MUTED, ha="left")

    # Inset.
    insets = spec.get("insets", [])
    panels = spec.get("panels", [])
    if insets:
        inset = insets[0]
        if inset.get("type") == "residual-basicblock":
            draw_basicblock(ax, 0.52, 0.42, inset.get("title", "Residual BasicBlock"), inset.get("subtitle", "repeated block"))
    elif any(panel.get("type") == "legend" for panel in panels):
        add_text(ax, 0.52, 1.62, "Legend", size=7.0, weight="bold", ha="left")
        add_text(ax, 0.52, 1.34, "feature stack: tensor scale\ncurved line: skip/fusion\nbadge: repeated block", size=5.5, color=MUTED, ha="left")

    add_text(ax, 5.35, 1.62, "Figure contract", size=7.0, weight="bold", ha="left")
    add_text(ax, 5.35, 1.35, spec.get("core_conclusion", ""), size=5.8, color=MUTED, ha="left")
    add_text(ax, 5.35, 0.95, "Caption-safe note", size=7.0, weight="bold", ha="left")
    caption_safe = audit.get("caption_safe_note") or "This diagram describes model topology and tensor-scale changes only; it does not claim accuracy."
    source_note = spec.get("source_note") or audit.get("source_paper") or audit.get("implementation_source") or ""
    add_text(ax, 5.35, 0.68, caption_safe, size=5.8, color=MUTED, ha="left")
    add_text(ax, 5.35, 0.24, source_note, size=5.3, color=MUTED, ha="left")

    out.parent.mkdir(parents=True, exist_ok=True)
    for ext in formats:
        path = out.with_suffix(f".{ext}")
        if ext == "png":
            fig.savefig(path, dpi=dpi, bbox_inches="tight", facecolor=PAPER)
        elif ext in {"svg", "pdf"}:
            fig.savefig(path, bbox_inches="tight", facecolor=PAPER)
        elif ext == "pptx":
            try:
                if pptx_mode == "native-shape":
                    write_native_pptx(spec, preset, path)
                else:
                    from pptx import Presentation  # type: ignore
                    from pptx.util import Inches  # type: ignore

                    png_path = out.with_suffix(".png")
                    if not png_path.exists():
                        fig.savefig(png_path, dpi=dpi, bbox_inches="tight", facecolor=PAPER)
                    prs = Presentation()
                    prs.slide_width = Inches(13.333)
                    prs.slide_height = Inches(7.5)
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    slide.shapes.add_picture(str(png_path), Inches(0.35), Inches(0.5), width=Inches(12.6))
                    prs.save(path)
            except Exception as exc:  # pragma: no cover - environment dependent
                path.with_suffix(".pptx.unavailable.txt").write_text(
                    f"PPTX export unavailable: {exc}\nInstall python-pptx or use the Presentations plugin for editable PPTX output.\n",
                    encoding="utf-8",
                )
    plt.close(fig)

    if qa_report:
        write_qa_report(spec, out, out.with_suffix(".qa.md"), formats, pptx_mode, explicit_unresolved)


def write_native_pptx(spec: dict[str, Any], preset: dict[str, Any], path: Path) -> None:
    from pptx import Presentation  # type: ignore
    from pptx.dml.color import RGBColor  # type: ignore
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE  # type: ignore
    from pptx.util import Inches, Pt  # type: ignore

    def rgb(hex_color: str) -> RGBColor:
        hex_color = hex_color.lstrip("#")
        return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    palette = preset["palette"]

    title = slide.shapes.add_textbox(Inches(0.35), Inches(0.25), Inches(8.5), Inches(0.35))
    title.text_frame.text = spec.get("title", spec.get("model", "Network architecture"))
    title.text_frame.paragraphs[0].font.size = Pt(18)
    title.text_frame.paragraphs[0].font.bold = True

    nodes = [spec.get("input", {"name": "Input", "shape": ""}), *spec.get("stages", []), *spec.get("head", [])]
    x0 = 0.55
    y0 = 2.45
    step = min(1.18, 11.6 / max(1, len(nodes)))
    prev_right = None
    for i, node in enumerate(nodes):
        x = Inches(x0 + i * step)
        y = Inches(y0 + min(i, 5) * 0.03)
        w = Inches(0.82)
        h = Inches(0.66)
        role = node.get("color_role", f"stage{i}") if isinstance(node, dict) else f"stage{i}"
        color = palette.get(role, "#E8EDF2")
        planes = min(int(node.get("planes", 3)) if isinstance(node, dict) else 3, 5)
        for p in range(planes):
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.035 * p), y - Inches(0.028 * p), w, h)
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb(color)
            shape.line.color.rgb = rgb(INK)
            shape.line.width = Pt(0.45)
        label = slide.shapes.add_textbox(x - Inches(0.08), y - Inches(0.42), Inches(1.05), Inches(0.28))
        label.text_frame.text = node.get("name", "node") if isinstance(node, dict) else "node"
        label.text_frame.paragraphs[0].font.size = Pt(7)
        label.text_frame.paragraphs[0].font.bold = True
        shape_text = slide.shapes.add_textbox(x - Inches(0.16), y + h + Inches(0.04), Inches(1.2), Inches(0.36))
        shape_text.text_frame.text = node.get("shape", "") if isinstance(node, dict) else ""
        shape_text.text_frame.paragraphs[0].font.size = Pt(6)
        if prev_right is not None:
            conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, prev_right, y + h / 2, x - Inches(0.08), y + h / 2)
            conn.line.color.rgb = rgb(INK)
            conn.line.width = Pt(0.8)
        prev_right = x + w + Inches(0.18)

    note = slide.shapes.add_textbox(Inches(0.55), Inches(5.8), Inches(11.8), Inches(0.55))
    note.text_frame.text = spec.get("caption", "Architecture schematic generated from .network.json.")
    note.text_frame.paragraphs[0].font.size = Pt(8)
    prs.save(path)


def write_qa_report(spec: dict[str, Any], out_base: Path, path: Path, formats: list[str], pptx_mode: str, unresolved_connections: list[str]) -> None:
    audit = spec.get("audit", {})
    output_checks = [(f"output_{ext}_exists", out_base.with_suffix(f".{ext}").exists()) for ext in formats]
    checks = [
        ("architecture_source_recorded", bool(spec.get("source_note") or audit.get("source_paper") or audit.get("implementation_source"))),
        ("core_conclusion_recorded", bool(spec.get("core_conclusion"))),
        ("input_shape_recorded", bool(spec.get("input", {}).get("shape"))),
        ("stage_shapes_recorded", all(stage.get("shape") for stage in spec.get("stages", []))),
        ("stage_channels_recorded", all("x" in stage.get("shape", "") for stage in spec.get("stages", []))),
        ("downsampling_marked", any(stage.get("downsample") for stage in spec.get("stages", []))),
        ("declared_connections_resolved", not unresolved_connections),
        ("inset_present_for_repeated_block", bool(spec.get("insets"))),
        ("caption_present", bool(spec.get("caption"))),
        ("caption_safe_note_recorded", bool(audit.get("caption_safe_note") or spec.get("caption"))),
        ("vector_export_requested", "svg" in formats or "pdf" in formats),
        ("pptx_mode_recorded", "pptx" not in formats or pptx_mode in {"image-backed", "native-shape"}),
        *output_checks,
    ]
    lines = [
        "# Network Architecture Figure QA",
        "",
        f"- Model: {spec.get('model', 'unknown')}",
        f"- Layout: {spec.get('layout', 'linear')}",
        f"- Requested formats: {', '.join(formats)}",
        f"- PPTX mode: {pptx_mode if 'pptx' in formats else 'not requested'}",
        f"- Unresolved connections: {', '.join(unresolved_connections) if unresolved_connections else 'none'}",
        "",
        "| Check | Status |",
        "|---|---|",
    ]
    for name, ok in checks:
        lines.append(f"| {name} | {'pass' if ok else 'review'} |")
    lines.extend(["", "## Caption", "", spec.get("caption", "Caption not provided.")])
    path.write_text("\n".join(lines) + "\n", encoding="utf-8")


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Render a neural-network architecture figure from JSON.")
    parser.add_argument("--spec", required=True, type=Path, help="Path to a .network.json structure spec.")
    parser.add_argument("--preset", default="ppt-template-rich", help="Preset name: thesis-clean, nature-minimal, ppt-template-rich.")
    parser.add_argument("--presets", type=Path, default=Path(__file__).resolve().parents[1] / "presets" / "network-figure-presets.json")
    parser.add_argument("--out", required=True, type=Path, help="Output path without extension or with any extension.")
    parser.add_argument("--formats", default="svg,pdf,png", help="Comma-separated formats: svg,pdf,png,pptx.")
    parser.add_argument("--pptx-mode", default="image-backed", choices=["image-backed", "native-shape"], help="PPTX export mode. image-backed embeds the PNG preview; native-shape creates editable PowerPoint shapes.")
    parser.add_argument("--qa-report", action="store_true", help="Write a Markdown QA report next to outputs.")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    spec = load_json(args.spec)
    presets = load_json(args.presets)
    if args.preset not in presets:
        raise SystemExit(f"Unknown preset {args.preset!r}; available: {', '.join(presets)}")
    formats = [f.strip().lower() for f in args.formats.split(",") if f.strip()]
    out = args.out.with_suffix("")
    render(spec, presets[args.preset], out, formats, args.qa_report, args.pptx_mode)


if __name__ == "__main__":
    main()
